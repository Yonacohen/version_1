from pptx import Presentation
def main():
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.tetx = "Hi"
    subtitle.text = "hi man"
    prs.save("test.pptx")


if __name__ == "__main__":
    main()



